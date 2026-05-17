from pptx import Presentation


class SlideAgent:
    def generate(self, topic: str, script: str, output_path: str):
        prs = Presentation()

        paragraphs = [p.strip() for p in script.split("\n") if p.strip()]

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic
        title_slide.placeholders[1].text = "AI Generated Educational Video"

        chunk_size = 4

        for i in range(0, len(paragraphs), chunk_size):
            chunk = paragraphs[i:i + chunk_size]

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Section {i // chunk_size + 1}"

            body = slide.placeholders[1]
            body.text = "\n".join(chunk)

        prs.save(output_path)

        return output_path