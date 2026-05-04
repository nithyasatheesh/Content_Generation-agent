import os
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Dict, Optional

import fitz  # PyMuPDF
from docx import Document
from openai import OpenAI
from pptx import Presentation
import streamlit as st


@dataclass
class GeneratedArtifacts:
    structured_content: str
    quiz_content: str
    slide_path: str
    video_path: Optional[str] = None


class ParserAgent:
    def detect_file_type(self, filename: str) -> Optional[str]:
        extension = os.path.splitext(filename)[1].lower()
        return {
            ".pdf": "pdf",
            ".docx": "docx",
            ".txt": "txt",
            ".pptx": "pptx",
        }.get(extension)

    def parse(self, file_type: str, path: str) -> str:
        parsers: Dict[str, Callable[[str], str]] = {
            "pdf": self._parse_pdf,
            "docx": self._parse_docx,
            "txt": self._parse_txt,
            "pptx": self._parse_pptx,
        }
        parser = parsers.get(file_type)
        if not parser:
            raise ValueError(f"Unsupported file type: {file_type}")
        return parser(path)

    def _parse_pdf(self, path: str) -> str:
        doc = fitz.open(path)
        return "".join(page.get_text() for page in doc)

    def _parse_docx(self, path: str) -> str:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    def _parse_txt(self, path: str) -> str:
        with open(path, "r", encoding="utf-8") as file:
            return file.read()

    def _parse_pptx(self, path: str) -> str:
        prs = Presentation(path)
        text_blocks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_blocks.append(shape.text)
        return "\n".join(text_blocks)


class ContentStructuringAgent:
    def __init__(self, client: OpenAI):
        self.client = client

    def structure(self, text: str) -> str:
        prompt = f"""
        Convert the following content into teaching slides.

        Rules:
        - Max 5 bullet points per slide
        - Each point under 12 words
        - Keep it simple and structured

        Content:
        {text[:4000]}
        """

        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
        )
        return response.choices[0].message.content or ""


class QuizGeneratorAgent:
    def __init__(self, client: OpenAI):
        self.client = client

    def generate(self, text: str) -> str:
        prompt = f"""
        Create a quiz:
        - 5 MCQs
        - Include answers

        Content:
        {text[:2000]}
        """

        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
        )
        return response.choices[0].message.content or ""


class SlideGeneratorAgent:
    def generate(self, structured_text: str, quiz_text: str, output_path: str) -> str:
        prs = Presentation()
        slides = [block for block in structured_text.split("\n\n") if block.strip()]

        for slide_text in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            lines = [line for line in slide_text.split("\n") if line.strip()]

            if lines:
                slide.shapes.title.text = lines[0]
                slide.placeholders[1].text = "\n".join(lines[1:])

        quiz_slide = prs.slides.add_slide(prs.slide_layouts[1])
        quiz_slide.shapes.title.text = "Quiz"
        quiz_slide.placeholders[1].text = quiz_text

        prs.save(output_path)
        return output_path


class VideoGeneratorAgent:
    def generate(self, structured_text: str, output_path: str) -> str:
        with open(output_path, "w", encoding="utf-8") as file:
            file.write("Video generation placeholder\n")
            file.write("Input summary:\n")
            file.write(structured_text[:500])
        return output_path


class ContentGenerationOrchestrator:
    def __init__(self, client: OpenAI) -> None:
        self.parser_agent = ParserAgent()
        self.structuring_agent = ContentStructuringAgent(client)
        self.slide_agent = SlideGeneratorAgent()
        self.quiz_agent = QuizGeneratorAgent(client)
        self.video_agent = VideoGeneratorAgent()

    def process(self, filename: str, temp_path: str, output_dir: str) -> GeneratedArtifacts:
        file_type = self.parser_agent.detect_file_type(filename)
        if not file_type:
            raise ValueError("Unsupported file type")

        raw_text = self.parser_agent.parse(file_type, temp_path)
        structured_content = self.structuring_agent.structure(raw_text)
        quiz_content = self.quiz_agent.generate(raw_text)

        slide_path = self.slide_agent.generate(
            structured_content,
            quiz_content,
            output_path=os.path.join(output_dir, "output.pptx"),
        )
        video_path = self.video_agent.generate(
            structured_content,
            output_path=os.path.join(output_dir, "output_video.txt"),
        )

        return GeneratedArtifacts(
            structured_content=structured_content,
            quiz_content=quiz_content,
            slide_path=slide_path,
            video_path=video_path,
        )


def run_app() -> None:
    st.set_page_config(page_title="Content Generation Agent", layout="wide")
    st.title("📚 Content Generation Agent")
    st.write("Upload PDF/DOCX/TXT/PPTX and generate structured content, quiz, and slides.")

    api_key = st.text_input("OpenAI API Key", type="password")
    uploaded_file = st.file_uploader("Upload source file", type=["pdf", "docx", "txt", "pptx"])

    if st.button("Generate Content", type="primary"):
        if not api_key:
            st.error("Please enter your OpenAI API key.")
            return
        if not uploaded_file:
            st.error("Please upload a file first.")
            return

        client = OpenAI(api_key=api_key)
        orchestrator = ContentGenerationOrchestrator(client)

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_input_path = os.path.join(temp_dir, uploaded_file.name)
            with open(temp_input_path, "wb") as temp_file:
                temp_file.write(uploaded_file.read())

            try:
                artifacts = orchestrator.process(uploaded_file.name, temp_input_path, temp_dir)
            except ValueError as error:
                st.error(str(error))
                return
            except Exception as error:  # noqa: BLE001
                st.exception(error)
                return

            st.success("Content package created.")

            col1, col2 = st.columns(2)
            with col1:
                st.subheader("Structured Content")
                st.text_area("", artifacts.structured_content, height=300)
            with col2:
                st.subheader("Quiz")
                st.text_area(" ", artifacts.quiz_content, height=300)

            with open(artifacts.slide_path, "rb") as slide_file:
                st.download_button(
                    "Download Slides (.pptx)",
                    data=slide_file,
                    file_name=Path(artifacts.slide_path).name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

            with open(artifacts.video_path, "rb") as video_file:
                st.download_button(
                    "Download Video Placeholder (.txt)",
                    data=video_file,
                    file_name=Path(artifacts.video_path).name,
                    mime="text/plain",
                )


if __name__ == "__main__":
    run_app()
